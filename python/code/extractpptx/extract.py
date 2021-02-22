#################################################
# Extract text from PowerPoint PPTX files
#
#
#
# You need to have the module "python-pptx" installed
# pip install python-pptx
#
# Load up the relevant modules
from pptx import Presentation
import glob
import sys

# Write the output of the extract to a text file
sys.stdout=open("output.txt","w")

# Extract the data from the file
for eachfile in glob.glob("*.pptx"):
    prs = Presentation(eachfile)
    print(eachfile)
    print("----------------------")
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(shape.text)

# Close the text file
sys.stdout.close()
